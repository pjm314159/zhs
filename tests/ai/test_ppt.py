"""Task 6.2 — ai/ppt.py TDD（python-pptx 本地提取）"""

from pathlib import Path
from unittest.mock import MagicMock, patch

from zhs.ai.ppt import PptConverter


class TestPptConverterInit:
    """PptConverter 初始化"""

    def test_default_params(self) -> None:
        """默认参数：本地模式，cleanup_local=True"""
        converter = PptConverter()
        assert converter._should_cleanup_local is True

    def test_cleanup_local_disabled(self) -> None:
        """cleanup_local=False"""
        converter = PptConverter(cleanup_local=False)
        assert converter._should_cleanup_local is False


class TestCleanupLocal:
    """本地文件清理"""

    def test_cleanup_local_deletes_file(self, tmp_path: Path) -> None:
        """cleanup_local=True 删除临时文件"""
        test_file = tmp_path / "test.pptx"
        test_file.write_bytes(b"fake ppt content")
        converter = PptConverter(cleanup_local=True)
        converter._cleanup_local(test_file)
        assert not test_file.exists()

    def test_cleanup_local_disabled(self, tmp_path: Path) -> None:
        """cleanup_local=False 时 convert 不调用 _cleanup_local"""
        test_file = tmp_path / "test.pptx"
        test_file.write_bytes(b"fake ppt content")
        converter = PptConverter(cleanup_local=False)
        with (
            patch.object(converter, "_download", return_value=test_file),
            patch.object(converter, "_extract_local", return_value="内容"),
            patch.object(converter, "_cleanup_local") as mock_cleanup,
        ):
            converter.convert("https://example.com/test.pptx")
            mock_cleanup.assert_not_called()

    def test_cleanup_nonexistent_file_no_error(self, tmp_path: Path) -> None:
        """清理不存在的文件不报错"""
        test_file = tmp_path / "nonexistent.pptx"
        converter = PptConverter(cleanup_local=True)
        converter._cleanup_local(test_file)  # 不应抛异常


class TestConvert:
    """完整转换流程"""

    @patch("zhs.ai.ppt.PptConverter._cleanup_local")
    @patch("zhs.ai.ppt.PptConverter._extract_local")
    @patch("zhs.ai.ppt.PptConverter._download")
    def test_convert_full_flow(
        self,
        mock_download: MagicMock,
        mock_extract: MagicMock,
        mock_cleanup: MagicMock,
    ) -> None:
        """完整流程：download → extract_local → cleanup_local"""
        mock_download.return_value = Path("/tmp/test.pptx")
        mock_extract.return_value = "PPT 文本内容"

        converter = PptConverter()
        result = converter.convert("https://example.com/test.pptx")

        assert result == "PPT 文本内容"
        mock_download.assert_called_once_with("https://example.com/test.pptx")
        mock_extract.assert_called_once_with(Path("/tmp/test.pptx"))
        mock_cleanup.assert_called_once_with(Path("/tmp/test.pptx"))

    @patch("zhs.ai.ppt.PptConverter._cleanup_local")
    @patch("zhs.ai.ppt.PptConverter._extract_local")
    @patch("zhs.ai.ppt.PptConverter._download")
    def test_convert_download_error_returns_empty(
        self,
        mock_download: MagicMock,
        mock_extract: MagicMock,
        mock_cleanup: MagicMock,
    ) -> None:
        """下载失败返回空字符串"""
        mock_download.side_effect = Exception("Network error")

        converter = PptConverter()
        result = converter.convert("https://example.com/bad.pptx")

        assert result == ""
        mock_extract.assert_not_called()


class TestExtractLocal:
    """python-pptx 本地提取"""

    @staticmethod
    def _create_pptx(path: Path, slides_data: list[list[str]]) -> None:
        """创建测试用 .pptx 文件"""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        for texts in slides_data:
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # 使用标题+内容布局
            for i, text in enumerate(texts):
                if i == 0 and slide.shapes.title:
                    slide.shapes.title.text = text
                elif len(slide.shapes.placeholders) > 1:
                    ph = slide.shapes.placeholders[1]
                    ph.text = text
                else:
                    from pptx.util import Emu

                    txBox = slide.shapes.add_textbox(Emu(0), Emu(0), Inches(5), Inches(1))
                    txBox.text_frame.text = text
        prs.save(str(path))

    def test_extract_local_single_slide(self, tmp_path: Path) -> None:
        """单张幻灯片提取"""
        pptx_path = tmp_path / "test.pptx"
        self._create_pptx(pptx_path, [["标题1", "内容1"]])

        converter = PptConverter()
        result = converter._extract_local(pptx_path)
        assert "标题1" in result
        assert "内容1" in result
        assert "[幻灯片 1]" in result

    def test_extract_local_multiple_slides(self, tmp_path: Path) -> None:
        """多张幻灯片提取"""
        pptx_path = tmp_path / "multi.pptx"
        self._create_pptx(pptx_path, [["标题1", "内容1"], ["标题2", "内容2"]])

        converter = PptConverter()
        result = converter._extract_local(pptx_path)
        assert "[幻灯片 1]" in result
        assert "[幻灯片 2]" in result
        assert "标题1" in result
        assert "标题2" in result

    def test_extract_local_empty_slide_skipped(self, tmp_path: Path) -> None:
        """空幻灯片不出现在结果中"""
        pptx_path = tmp_path / "empty.pptx"
        self._create_pptx(pptx_path, [["有内容"], []])

        converter = PptConverter()
        result = converter._extract_local(pptx_path)
        assert "[幻灯片 1]" in result
        assert "[幻灯片 2]" not in result

    def test_extract_local_table(self, tmp_path: Path) -> None:
        """表格内容提取"""
        from pptx import Presentation
        from pptx.util import Inches

        pptx_path = tmp_path / "table.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # 空白布局
        rows, cols = 2, 2
        table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(4), Inches(2))
        table = table_shape.table
        table.cell(0, 0).text = "姓名"
        table.cell(0, 1).text = "分数"
        table.cell(1, 0).text = "张三"
        table.cell(1, 1).text = "95"
        prs.save(str(pptx_path))

        converter = PptConverter()
        result = converter._extract_local(pptx_path)
        assert "姓名" in result
        assert "张三" in result
        assert "95" in result

    def test_convert_local_mode(self, tmp_path: Path) -> None:
        """本地模式完整转换流程"""
        pptx_path = tmp_path / "local.pptx"
        self._create_pptx(pptx_path, [["测试标题", "测试内容"]])

        converter = PptConverter(cleanup_local=False)
        with patch.object(converter, "_download", return_value=pptx_path):
            result = converter.convert("https://example.com/test.pptx")

        assert "测试标题" in result
        assert "测试内容" in result
