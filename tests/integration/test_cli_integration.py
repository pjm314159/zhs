"""CLI 集成测试 — typer 命令行端到端"""

import json
from pathlib import Path

import pytest
from typer.testing import CliRunner

from zhs.__main__ import app

pytestmark = pytest.mark.integration

runner = CliRunner()


class TestCLIFetch:
    """课程列表获取"""

    def test_fetch_course_list(self) -> None:
        """C-02: zhs -f 获取课程列表"""
        result = runner.invoke(app, ["-f"])
        # 可能输出包含课程信息或错误（未登录）
        assert result.exit_code == 0 or "登录" in result.output or "Cookie" in result.output

    def test_fetch_execution_json_format(self) -> None:
        """C-03: execution.json 保存为三类 dict"""
        from zhs.utils.path import get_data_dir

        exec_path = get_data_dir() / "execution.json"
        if not exec_path.exists():
            pytest.skip("先运行 zhs -f 生成 execution.json")

        with open(exec_path, encoding="utf-8") as f:
            data = json.load(f)

        assert isinstance(data, dict)
        assert "zhidao" in data
        assert "hike" in data
        assert "ai" in data
        assert isinstance(data["zhidao"], list)
        assert isinstance(data["hike"], list)
        assert isinstance(data["ai"], list)


class TestCLILogin:
    """登录命令"""

    def test_login_command_exists(self) -> None:
        """C-01: zhs login 子命令存在"""
        result = runner.invoke(app, ["login", "--help"])
        assert result.exit_code == 0
        assert "login" in result.output.lower() or "扫码" in result.output


class TestCLINotLoggedIn:
    """未登录状态"""

    def test_not_logged_in_shows_error(self) -> None:
        """C-11: 未登录时提示需要登录（通过 mock 验证）"""
        from unittest.mock import patch

        with patch("zhs.__main__._try_restore_cookies", return_value=False):
            result = runner.invoke(app, [], catch_exceptions=True)
            output_ok = "登录" in result.output or "Cookie" in result.output or "login" in result.output.lower()
            assert output_ok or result.exit_code != 0


class TestCLICourseType:
    """课程类型检测"""

    def test_type_flag_in_help(self) -> None:
        """C-13/C-14: --type 参数在帮助中"""
        result = runner.invoke(app, ["--help"])
        assert result.exit_code == 0
        assert "--type" in result.output

    def test_ai_course_and_class_flags(self) -> None:
        """C-12: --ai-course 和 --ai-class 参数存在"""
        result = runner.invoke(app, ["--help"])
        assert "--ai-course" in result.output
        assert "--ai-class" in result.output

    def test_noexam_flag(self) -> None:
        """C-09: --noexam 参数存在"""
        result = runner.invoke(app, ["--help"])
        assert "--noexam" in result.output

    def test_no_ai_exam_flag(self) -> None:
        """C-10: --no-ai-exam 参数存在"""
        result = runner.invoke(app, ["--help"])
        assert "--no-ai-exam" in result.output


class TestCLIPptLocal:
    """PPT 本地转换（python-pptx）"""

    def test_ppt_converter_local_mode(self, tmp_path: Path) -> None:
        """python-pptx 本地模式无需 API Key"""
        from pptx import Presentation

        from zhs.ai.ppt import PptConverter

        # 创建测试 PPT
        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title:
            slide.shapes.title.text = "集成测试标题"
        prs.save(str(pptx_path))

        # 本地模式转换
        converter = PptConverter(cleanup_local=False)
        from unittest.mock import patch

        with patch.object(converter, "_download", return_value=pptx_path):
            result = converter.convert("https://example.com/test.pptx")

        assert "集成测试标题" in result
