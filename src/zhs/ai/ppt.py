"""PPT 转文本（本地 python-pptx）"""

from pathlib import Path

from loguru import logger

from zhs.utils.path import get_data_dir


class PptConverter:
    """PPT 转文本（使用 python-pptx 本地提取，无需 API Key）"""

    def __init__(
        self,
        cleanup_local: bool = True,
    ) -> None:
        self._should_cleanup_local = cleanup_local
        self._download_path = get_data_dir() / "AiDownloadCache"

    def convert(self, url: str) -> str:
        """下载 PPT 并转为文本，完成后清理本地临时文件"""
        try:
            local_path = self._download(url)
            text = self._extract_local(local_path)
            if self._should_cleanup_local:
                self._cleanup_local(local_path)
            return text
        except Exception as e:
            logger.error(f"PPT 转换失败: {e}")
            return ""

    def _download(self, url: str) -> Path:
        """下载 PPT 文件到本地"""
        from urllib.parse import urlparse

        import httpx

        parsed = urlparse(url)
        file_path = parsed.path.lstrip("/")
        local_path = self._download_path / file_path

        if local_path.exists():
            logger.info(f"文件已存在: {local_path}")
            return local_path

        local_path.parent.mkdir(parents=True, exist_ok=True)

        with httpx.stream("GET", url) as response:
            response.raise_for_status()
            with open(local_path, "wb") as f:
                for chunk in response.iter_bytes(chunk_size=8192):
                    f.write(chunk)

        logger.info(f"文件已下载: {local_path}")
        return local_path

    def _extract_local(self, file_path: Path) -> str:
        """使用 python-pptx 本地提取 PPT 文本"""
        from pptx import Presentation

        prs = Presentation(str(file_path))
        texts: list[str] = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_texts:
                            slide_texts.append(" | ".join(row_texts))
            if slide_texts:
                texts.append(f"[幻灯片 {slide_idx}]\n" + "\n".join(slide_texts))
        return "\n\n".join(texts)

    def _cleanup_local(self, file_path: Path) -> None:
        """清理本地临时文件"""
        try:
            if file_path.exists():
                file_path.unlink()
                logger.debug(f"本地文件已清理: {file_path}")
        except Exception as e:
            logger.error(f"清理本地文件失败: {e}")
